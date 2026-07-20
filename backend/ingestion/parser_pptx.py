"""PPTX → Markdown 友好文本（，python-pptx）。

- 只提取 has_text_frame 的文本 + 表格；
- 每页输出 「## 第N页 [标题]」；
- 不提取图片/SmartArt/图表/演讲者备注。
"""
from io import BytesIO

from pptx import Presentation


def parse_pptx(file_bytes: bytes, filename: str) -> tuple[str, str]:
    """返回 (文本, 标题=文件名)。"""
    try:
        prs = Presentation(BytesIO(file_bytes))
    except Exception as e:
        raise ValueError(f"无法解析 PPTX 文件: {e}") from e

    pages = []
    for idx, slide in enumerate(prs.slides, start=1):
        title = ""
        if slide.shapes.title is not None:
            title = slide.shapes.title.text.strip()

        parts = [f"## 第{idx}页" + (f" {title}" if title else "")]
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt:
                    parts.append(txt)
            elif shape.has_table:
                rows = []
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(" | ".join(cells))
                parts.append("\n".join(rows))
        pages.append("\n\n".join(parts))

    return "\n\n".join(pages), filename
