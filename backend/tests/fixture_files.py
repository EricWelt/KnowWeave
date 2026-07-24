"""测试文件生成器：用真实库合成 PDF/PPTX/MD 二进制内容。

放在 tests 里而非提交二进制文件 —— 保证仓库干净、文件可随时再生成。
"""
from io import BytesIO


def make_pdf_bytes() -> bytes:
    import fitz

    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Operating System Process Scheduling")
    page.insert_text((72, 100), "First Come First Served is a scheduling algorithm.")
    buf = BytesIO()
    doc.save(buf)
    doc.close()
    return buf.getvalue()


def make_pptx_bytes() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Deadlock Basics"
    slide.placeholders[1].text = "Four necessary conditions of deadlock"
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def make_markdown_bytes() -> bytes:
    content = "# 进程管理\n\n进程是资源分配的基本单位。\n\n- 进程控制块 PCB\n- 上下文切换"
    return content.encode("utf-8")


def make_garbage_bytes() -> bytes:
    return b"this is not a real file format at all " * 10
